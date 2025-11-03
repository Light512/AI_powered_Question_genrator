import json
import torch
import uvicorn
from fastapi import FastAPI, File, UploadFile
from pptx import Presentation
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline
from huggingface_hub import login

# ============================================
# Login to Hugging Face
# ============================================
login(token="yourrrr_tokkeen")

# ============================================
# Load Llama-3.1 Model
# ============================================
MODEL_ID = "meta-llama/Llama-3.1-8B-Instruct"

tokenizer = AutoTokenizer.from_pretrained(MODEL_ID)
model = AutoModelForCausalLM.from_pretrained(
    MODEL_ID,
    torch_dtype=torch.bfloat16,
    device_map="auto"
)

pipe = pipeline(
    "text-generation",
    model=model,
    tokenizer=tokenizer,
    max_new_tokens=512,
    temperature=0.3,
    top_p=0.9,
    repetition_penalty=1.1,
)

# ============================================
# FastAPI Application
# ============================================
app = FastAPI()


def extract_text_from_pptx(pptx_file):
    """Extracts text from a PPTX resume."""
    text = []
    ppt = Presentation(pptx_file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def generate_questions_llama(resume_text, job_description, experience_level="mid"):
    """Generates interview questions using the Llama model."""
    prompt = f"""
You are an AI interview assistant. Generate tailored interview questions based on the candidate’s resume and job description.

Resume: {resume_text}

Job Description: {job_description}

Experience Level: {experience_level}

Generate 5 technical and 3 behavioral interview questions in a clear list format.
"""

    messages = [
        {"role": "system", "content": "You are an AI assistant that generates interview questions."},
        {"role": "user", "content": prompt}
    ]

    response = pipe(messages)
    return response[0]["generated_text"]


@app.post("/generate-questions/")
async def generate_interview_questions(file: UploadFile = File(...), job_description: str = "Software Engineer"):
    """API endpoint to generate interview questions from a PPTX resume."""
    resume_text = extract_text_from_pptx(file.file)
    experience_level = "senior" if "10+ years" in resume_text else "mid"

    questions = generate_questions_llama(resume_text, job_description, experience_level)
    return {"questions": questions}
